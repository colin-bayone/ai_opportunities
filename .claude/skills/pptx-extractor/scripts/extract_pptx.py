#!/usr/bin/env python3
"""
PPTX Extractor - Faithful slide-by-slide markdown extraction using Gemini Vision.

Converts PowerPoint slides to images via LibreOffice, then sends each slide
in parallel to Gemini 2.5 Pro for faithful markdown recreation, ASCII layout,
and visual element descriptions.

Usage:
    python extract_pptx.py <file.pptx> [--slides 1-5,10] [--model gemini-2.5-pro] [--dry-run]
    python extract_pptx.py <file.pptx> --output-dir <custom_dir>  # override default output location

Environment:
    GEMINI_API_KEY - Required for Gemini API access (can be in .env file)

Dependencies:
    - python-pptx (speaker notes extraction)
    - google-genai (Gemini API)
    - libreoffice (headless, for PPTX to PDF conversion)
    - pdftoppm (poppler-utils, for PDF to PNG conversion)

Output (created next to the source PPTX by default):
    <pptx_stem>/
    ├── index.md
    ├── metadata.json
    ├── slide_01/
    │   ├── slide_01.png
    │   ├── content.md
    │   ├── layout.md
    │   └── visual_elements.md
    ├── slide_02/
    │   └── ...
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path


def load_env():
    """Load environment variables from .env file if available."""
    try:
        from dotenv import load_dotenv
        env_path = Path.cwd() / ".env"
        if not env_path.exists():
            for parent in Path.cwd().parents:
                candidate = parent / ".env"
                if candidate.exists():
                    env_path = candidate
                    break
        if env_path.exists():
            load_dotenv(env_path)
    except ImportError:
        pass


load_env()


def check_gemini_available() -> bool:
    """Check if Gemini SDK is available."""
    try:
        from google import genai
        return True
    except ImportError:
        return False


def get_gemini_client():
    """Initialize Gemini client."""
    from google import genai
    api_key = os.environ.get("GEMINI_API_KEY")
    if not api_key:
        raise ValueError("GEMINI_API_KEY environment variable not set")
    return genai.Client(api_key=api_key)


def check_system_deps() -> list[str]:
    """Check for required system dependencies. Returns list of missing deps."""
    missing = []
    if not shutil.which("libreoffice"):
        missing.append("libreoffice (apt install libreoffice-impress)")
    if not shutil.which("pdftoppm"):
        missing.append("pdftoppm (apt install poppler-utils)")
    return missing


def parse_slide_range(slide_spec: str, max_slides: int) -> list[int]:
    """Parse slide specification like '1-5,10,15-20' into list of slide numbers."""
    if not slide_spec:
        return list(range(1, max_slides + 1))

    slides = set()
    for part in slide_spec.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            start = int(start)
            end = min(int(end), max_slides)
            slides.update(range(start, end + 1))
        else:
            slide_num = int(part)
            if slide_num <= max_slides:
                slides.add(slide_num)

    return sorted(slides)


def get_pptx_info(file_path: Path) -> dict:
    """Extract structural info from PPTX using python-pptx."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(file_path))
    total_slides = len(prs.slides)

    slides_info = []
    for i, slide in enumerate(prs.slides):
        # Get title
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    title = t[:120]
                    break

        # Get speaker notes
        notes = None
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text

        # Count shapes and images
        shape_count = len(slide.shapes)
        image_count = sum(
            1 for s in slide.shapes
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        )

        # Check for tables and charts
        has_table = any(s.has_table for s in slide.shapes)
        has_chart = any(s.has_chart for s in slide.shapes)

        slides_info.append({
            "number": i + 1,
            "title": title,
            "speaker_notes": notes,
            "shape_count": shape_count,
            "image_count": image_count,
            "has_table": has_table,
            "has_chart": has_chart,
        })

    return {
        "total_slides": total_slides,
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides_info,
    }


def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    """Convert PPTX to PDF using LibreOffice headless."""
    result = subprocess.run(
        [
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(output_dir), str(pptx_path)
        ],
        capture_output=True, text=True, timeout=300
    )

    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

    pdf_name = pptx_path.stem + ".pdf"
    pdf_path = output_dir / pdf_name

    if not pdf_path.exists():
        raise RuntimeError(f"PDF not created at expected path: {pdf_path}")

    return pdf_path


def extract_slide_images(pdf_path: Path, output_dir: Path, slides: list[int], dpi: int = 300) -> dict[int, Path]:
    """Extract specific slides from PDF as PNG images using pdftoppm."""
    output_dir.mkdir(parents=True, exist_ok=True)
    slide_images = {}

    for slide_num in slides:
        prefix = output_dir / f"slide_{slide_num:02d}"
        result = subprocess.run(
            [
                "pdftoppm", "-png",
                "-f", str(slide_num), "-l", str(slide_num),
                "-r", str(dpi),
                str(pdf_path), str(prefix)
            ],
            capture_output=True, text=True, timeout=60
        )

        if result.returncode != 0:
            print(f"  Warning: pdftoppm failed for slide {slide_num}: {result.stderr}", file=sys.stderr)
            continue

        # pdftoppm names output as prefix-NN.png
        expected = prefix.parent / f"{prefix.name}-{slide_num:02d}.png"
        final = prefix.parent / f"slide_{slide_num:02d}.png"

        if expected.exists():
            expected.rename(final)
            slide_images[slide_num] = final
        else:
            # Try to find whatever pdftoppm created
            candidates = list(prefix.parent.glob(f"{prefix.name}*.png"))
            if candidates:
                candidates[0].rename(final)
                slide_images[slide_num] = final
            else:
                print(f"  Warning: No PNG created for slide {slide_num}", file=sys.stderr)

    return slide_images


def process_slide_with_gemini(
    client,
    image_path: Path,
    slide_num: int,
    model: str = "gemini-2.5-pro"
) -> dict:
    """Send a single slide image to Gemini for extraction."""
    from google.genai import types

    image_bytes = image_path.read_bytes()

    prompt = f"""You are analyzing slide {slide_num} of a PowerPoint presentation, rendered as an image.

Provide THREE clearly separated sections:

## Markdown Content
Recreate the complete visible content of this slide in markdown.
Include all text, bullet points, tables, and data exactly as shown.
Do NOT summarize. Capture every word faithfully.
Use proper markdown formatting (headers, lists, bold, etc.) to match the slide structure.

## ASCII Layout
Show the spatial arrangement of elements on this slide using ASCII art.
Indicate where titles, text blocks, images, and graphics are positioned relative to each other.
Use box-drawing characters and labels.
Show approximate proportions.

## Visual Elements
List and describe every image, screenshot, diagram, chart, icon, or graphic visible on this slide.
For each visual element, provide:
- Type (photo, icon, diagram, chart, screenshot, logo, etc.)
- Position on slide (top-left, center, bottom-right, etc.)
- Description of what it depicts
- If it contains text, transcribe that text

Be thorough and faithful. This is for document archival, not summarization."""

    try:
        response = client.models.generate_content(
            model=model,
            contents=[
                types.Part.from_bytes(data=image_bytes, mime_type="image/png"),
                prompt
            ]
        )
        return {"success": True, "text": response.text, "slide_num": slide_num}

    except Exception as e:
        return {"success": False, "text": f"Error processing slide {slide_num}: {str(e)}", "slide_num": slide_num}


def parse_gemini_response(text: str) -> dict:
    """Parse Gemini response into markdown, layout, and visual elements sections.

    Only splits on our three known section headers, not arbitrary ## headers
    that Gemini may use inside the markdown content itself.
    """
    sections = {"markdown": "", "layout": "", "visual_elements": ""}

    # Only match our specific section headers, not any ## in the content
    section_pattern = re.compile(
        r'^## (?:Markdown Content|ASCII Layout|Visual Elements)\s*$',
        re.MULTILINE
    )
    matches = list(section_pattern.finditer(text))

    for i, match in enumerate(matches):
        header = match.group(0).strip().lower()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        content = text[start:end].strip()

        if "markdown" in header:
            sections["markdown"] = content
        elif "ascii" in header or "layout" in header:
            sections["layout"] = content
        elif "visual" in header:
            sections["visual_elements"] = content

    # If parsing failed, put everything in markdown
    if not any(sections.values()):
        sections["markdown"] = text

    return sections


def write_outputs(
    output_dir: Path,
    pptx_name: str,
    pptx_info: dict,
    processed_slides: dict,
    slide_images: dict[int, Path],
    model: str,
) -> dict:
    """Write all output files: one folder per slide, each with its own files."""
    output_dir.mkdir(parents=True, exist_ok=True)
    files_created = []

    for slide_num in sorted(processed_slides.keys()):
        result = processed_slides[slide_num]
        slide_info = pptx_info["slides"][slide_num - 1]
        sections = parse_gemini_response(result["text"])

        # Create per-slide folder
        slide_dir = output_dir / f"slide_{slide_num:02d}"
        slide_dir.mkdir(parents=True, exist_ok=True)

        # Copy slide image into its folder
        if slide_num in slide_images:
            dest = slide_dir / f"slide_{slide_num:02d}.png"
            shutil.copy2(slide_images[slide_num], dest)
            files_created.append(str(dest))

        # raw.md - always save the complete Gemini response, no parsing
        raw_path = slide_dir / "raw.md"
        raw_path.write_text(result["text"])
        files_created.append(str(raw_path))

        # content.md - faithful markdown recreation
        content_path = slide_dir / "content.md"
        content = f"# Slide {slide_num}"
        if slide_info["title"]:
            content += f" - {slide_info['title'][:80]}"
        content += "\n\n"
        if slide_info.get("speaker_notes"):
            content += f"> **Speaker Notes:** {slide_info['speaker_notes']}\n\n---\n\n"
        content += sections["markdown"]
        content += f"\n\n---\n*Source: {pptx_name}, Slide {slide_num} of {pptx_info['total_slides']}*\n"
        content_path.write_text(content)
        files_created.append(str(content_path))

        # layout.md - ASCII spatial layout
        layout_path = slide_dir / "layout.md"
        layout = f"# Slide {slide_num} - Spatial Layout\n\n"
        layout += sections["layout"]
        layout += "\n"
        layout_path.write_text(layout)
        files_created.append(str(layout_path))

        # visual_elements.md - image/graphic descriptions
        visual_path = slide_dir / "visual_elements.md"
        visual = f"# Slide {slide_num} - Visual Elements\n\n"
        visual += sections["visual_elements"]
        visual += "\n"
        visual_path.write_text(visual)
        files_created.append(str(visual_path))

    # Write metadata.json
    metadata = {
        "source_file": pptx_name,
        "total_slides": pptx_info["total_slides"],
        "slides_processed": len(processed_slides),
        "model": model,
        "slides": [],
    }

    for slide_num in sorted(processed_slides.keys()):
        result = processed_slides[slide_num]
        slide_info = pptx_info["slides"][slide_num - 1]

        metadata["slides"].append({
            "number": slide_num,
            "title": slide_info["title"],
            "speaker_notes": slide_info.get("speaker_notes"),
            "shape_count": slide_info["shape_count"],
            "image_count": slide_info["image_count"],
            "has_table": slide_info["has_table"],
            "has_chart": slide_info["has_chart"],
            "extraction_success": result["success"],
        })

    metadata_path = output_dir / "metadata.json"
    metadata_path.write_text(json.dumps(metadata, indent=2))
    files_created.append(str(metadata_path))

    # Write index.md
    index_lines = [
        f"# {pptx_name}\n",
        f"**Total Slides:** {pptx_info['total_slides']}",
        f"**Slides Extracted:** {len(processed_slides)}",
        f"**Model:** {model}\n",
        "## Slides\n",
    ]

    for slide_num in sorted(processed_slides.keys()):
        slide_info = pptx_info["slides"][slide_num - 1]
        title = slide_info["title"][:60] if slide_info["title"] else "(untitled)"
        notes_marker = " [has notes]" if slide_info.get("speaker_notes") else ""
        folder = f"slide_{slide_num:02d}"
        index_lines.append(
            f"- **Slide {slide_num:02d}** - {title}{notes_marker}\n"
            f"  - [Content](./{folder}/content.md) "
            f"| [Layout](./{folder}/layout.md) "
            f"| [Visual Elements](./{folder}/visual_elements.md) "
            f"| [Image](./{folder}/slide_{slide_num:02d}.png)"
        )

    index_lines.append(f"\n---\n*Extracted by PPTX Extractor*\n")

    index_path = output_dir / "index.md"
    index_path.write_text("\n".join(index_lines))
    files_created.insert(0, str(index_path))

    return {"files_created": files_created}


def main():
    parser = argparse.ArgumentParser(
        description="Extract PowerPoint slides to markdown using Gemini Vision"
    )
    parser.add_argument("file_path", help="Path to PPTX or PDF file")
    parser.add_argument("--output-dir", "-o", default=None,
                        help="Output directory (default: <pptx_name>/ next to source file)")
    parser.add_argument("--slides", "-s", default=None,
                        help="Slide range to process (e.g., '1-5,10')")
    parser.add_argument("--model", "-m", default="gemini-2.5-pro",
                        help="Gemini model to use (default: gemini-2.5-pro)")
    parser.add_argument("--dpi", type=int, default=300,
                        help="DPI for slide image rendering (default: 300)")
    parser.add_argument("--workers", "-w", type=int, default=5,
                        help="Max parallel Gemini API calls (default: 5)")
    parser.add_argument("--dry-run", action="store_true",
                        help="Show what would be processed without calling API")

    args = parser.parse_args()

    file_path = Path(args.file_path).resolve()

    # Default output dir: same parent as PPTX, folder named after the file
    if args.output_dir:
        output_dir = Path(args.output_dir)
    else:
        output_dir = file_path.parent / file_path.stem

    # Validate input
    if not file_path.exists():
        print(json.dumps({"error": f"File not found: {file_path}"}))
        sys.exit(1)

    is_pdf = file_path.suffix.lower() == ".pdf"
    is_pptx = file_path.suffix.lower() == ".pptx"

    if not is_pdf and not is_pptx:
        print(json.dumps({"error": f"Unsupported file type: {file_path.suffix}", "supported": [".pptx", ".pdf"]}))
        sys.exit(1)

    # Check system dependencies (LibreOffice only needed for PPTX)
    missing = check_system_deps()
    if is_pdf:
        missing = [m for m in missing if "libreoffice" not in m]
    if missing and not args.dry_run:
        print(json.dumps({"error": "Missing system dependencies", "missing": missing}))
        sys.exit(1)

    # Check Gemini availability
    if not args.dry_run and not check_gemini_available():
        print(json.dumps({"error": "google-genai not installed", "fix": "pip install google-genai"}))
        sys.exit(1)

    # Get slide/page info
    print(f"Reading {file_path.name}...", file=sys.stderr)

    if is_pptx:
        pptx_info = get_pptx_info(file_path)
    else:
        # PDF mode: count pages, no speaker notes or shape metadata
        page_count_bytes = file_path.read_bytes()
        total_pages = page_count_bytes.count(b"/Type/Page") + page_count_bytes.count(b"/Type /Page")
        total_pages = max(total_pages, 1)
        pptx_info = {
            "total_slides": total_pages,
            "slide_width": None,
            "slide_height": None,
            "slides": [
                {
                    "number": i + 1,
                    "title": "",
                    "speaker_notes": None,
                    "shape_count": 0,
                    "image_count": 0,
                    "has_table": False,
                    "has_chart": False,
                }
                for i in range(total_pages)
            ],
        }

    total_slides = pptx_info["total_slides"]

    # Parse slide range
    slides_to_process = parse_slide_range(args.slides, total_slides)

    # Dry run
    if args.dry_run:
        print(json.dumps({
            "dry_run": True,
            "file": str(file_path),
            "input_type": "pdf" if is_pdf else "pptx",
            "total_slides": total_slides,
            "slides_to_process": slides_to_process,
            "output_dir": str(output_dir),
            "model": args.model,
            "dpi": args.dpi,
            "workers": args.workers,
            "slides": [
                {
                    "number": info["number"],
                    "title": info["title"][:60],
                    "has_notes": info["speaker_notes"] is not None,
                    "shapes": info["shape_count"],
                    "images": info["image_count"],
                }
                for info in pptx_info["slides"]
                if info["number"] in slides_to_process
            ]
        }, indent=2))
        sys.exit(0)

    # Create temp directory for intermediate files
    with tempfile.TemporaryDirectory(prefix="pptx_extract_") as temp_dir:
        temp_path = Path(temp_dir)

        # Step 1: Get a PDF (convert if PPTX, use directly if PDF)
        if is_pptx:
            print(f"Converting to PDF via LibreOffice...", file=sys.stderr)
            pdf_path = convert_pptx_to_pdf(file_path, temp_path)
            print(f"  PDF created: {pdf_path.stat().st_size / 1024 / 1024:.1f} MB", file=sys.stderr)
        else:
            pdf_path = file_path
            print(f"  Using PDF directly: {pdf_path.stat().st_size / 1024 / 1024:.1f} MB", file=sys.stderr)

        # Step 2: Extract slide images
        print(f"Extracting {len(slides_to_process)} slide images at {args.dpi} DPI...", file=sys.stderr)
        slide_images = extract_slide_images(pdf_path, temp_path / "images", slides_to_process, args.dpi)
        print(f"  {len(slide_images)} images extracted", file=sys.stderr)

        if not slide_images:
            print(json.dumps({"error": "No slide images were extracted"}))
            sys.exit(1)

        # Step 3: Process slides with Gemini in parallel
        client = get_gemini_client()
        processed_slides = {}
        total = len(slide_images)
        completed = 0

        print(f"Processing {total} slides with Gemini ({args.workers} parallel workers)...", file=sys.stderr)

        with ThreadPoolExecutor(max_workers=args.workers) as executor:
            futures = {
                executor.submit(
                    process_slide_with_gemini, client, slide_images[slide_num], slide_num, args.model
                ): slide_num
                for slide_num in sorted(slide_images.keys())
            }

            for future in as_completed(futures):
                slide_num = futures[future]
                slide_info = pptx_info["slides"][slide_num - 1]
                title_preview = slide_info["title"][:50] if slide_info["title"] else "(untitled)"

                try:
                    result = future.result()
                    processed_slides[slide_num] = result
                    completed += 1
                    status = "ok" if result["success"] else "FAILED"
                    print(
                        f"  [{completed}/{total}] Slide {slide_num}: {title_preview} - {status} ({len(result['text'])} chars)",
                        file=sys.stderr, flush=True
                    )
                except Exception as e:
                    completed += 1
                    processed_slides[slide_num] = {
                        "success": False,
                        "text": f"Error: {str(e)}",
                        "slide_num": slide_num,
                    }
                    print(
                        f"  [{completed}/{total}] Slide {slide_num}: {title_preview} - EXCEPTION: {e}",
                        file=sys.stderr, flush=True
                    )

        # Step 4: Write organized output
        print(f"Writing output to {output_dir}...", file=sys.stderr)
        output_info = write_outputs(
            output_dir, file_path.name, pptx_info,
            processed_slides, slide_images, args.model
        )

    # Summary
    successful = sum(1 for r in processed_slides.values() if r["success"])
    failed = len(processed_slides) - successful

    result = {
        "success": True,
        "file": str(file_path),
        "total_slides": total_slides,
        "slides_processed": len(processed_slides),
        "slides_successful": successful,
        "slides_failed": failed,
        "output_dir": str(output_dir),
        "files_created": output_info["files_created"],
    }

    print(json.dumps(result, indent=2))
    print(f"\nDone! {successful}/{len(processed_slides)} slides extracted successfully.", file=sys.stderr)


if __name__ == "__main__":
    main()
